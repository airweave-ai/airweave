"""PowerPoint bongo implementation.

Creates, updates, and deletes test PowerPoint presentations via the Microsoft Graph API.
"""

import asyncio
import io
import time
import uuid
from typing import Any, Dict, List

import httpx
from monke.bongos.base_bongo import BaseBongo
from monke.generation.powerpoint import generate_presentations_content
from monke.utils.logging import get_logger

# Try to import python-pptx for PowerPoint presentation creation
try:
    from pptx import Presentation
    from pptx.util import Inches

    HAS_PYTHON_PPTX = True
except ImportError:
    HAS_PYTHON_PPTX = False

GRAPH = "https://graph.microsoft.com/v1.0"


class PowerPointBongo(BaseBongo):
    """Bongo for PowerPoint that creates test entities for E2E testing.

    Key responsibilities:
    - Create test PowerPoint presentations in OneDrive
    - Add slides with content containing verification tokens
    - Update presentations to test incremental sync
    - Delete presentations to test deletion detection
    - Clean up all test data
    """

    connector_type = "powerpoint"

    def __init__(self, credentials: Dict[str, Any], **kwargs):
        super().__init__(credentials)
        self.access_token: str = credentials["access_token"]
        self.entity_count: int = int(kwargs.get("entity_count", 3))
        self.openai_model: str = kwargs.get("openai_model", "gpt-4.1-mini")
        self.rate_limit_delay = float(kwargs.get("rate_limit_delay_ms", 500)) / 1000.0
        self.logger = get_logger("powerpoint_bongo")

        # Track created resources for cleanup
        self._test_presentations: List[Dict[str, Any]] = []
        self._last_req = 0.0

        if not HAS_PYTHON_PPTX:
            raise ImportError(
                "python-pptx is required for PowerPoint bongo. Install with: pip install python-pptx"
            )

    async def create_entities(self) -> List[Dict[str, Any]]:
        """Create test PowerPoint presentations in OneDrive."""
        self.logger.info(
            f"🥁 Creating {self.entity_count} PowerPoint test presentations"
        )
        out: List[Dict[str, Any]] = []

        # Generate tokens for each presentation
        tokens = [uuid.uuid4().hex[:8] for _ in range(self.entity_count)]

        # Generate presentation content
        test_name = f"TestPresentation_{uuid.uuid4().hex[:8]}"
        filenames, presentation_content = await generate_presentations_content(
            self.openai_model, tokens, test_name
        )

        self.logger.info(f"📊 Generated {len(presentation_content)} presentations")

        async with httpx.AsyncClient(base_url=GRAPH, timeout=60) as client:
            for i, (filename, pres_content, token) in enumerate(
                zip(filenames, presentation_content, tokens)
            ):
                await self._pace()

                # Sanitize filename for OneDrive (remove illegal characters)
                safe_filename = self._sanitize_filename(filename)

                # Create PowerPoint presentation file
                pres_bytes = self._create_powerpoint_file(pres_content)

                # Upload to OneDrive
                self.logger.info(
                    f"📤 Uploading PowerPoint presentation: {safe_filename}"
                )
                upload_url = f"/me/drive/root:/{safe_filename}:/content"

                r = await client.put(
                    upload_url,
                    headers={
                        "Authorization": f"Bearer {self.access_token}",
                        "Content-Type": (
                            "application/vnd.openxmlformats-officedocument."
                            "presentationml.presentation"
                        ),
                    },
                    content=pres_bytes,
                )

                if r.status_code not in (200, 201):
                    self.logger.error(f"Upload failed {r.status_code}: {r.text}")
                    r.raise_for_status()

                pres_file = r.json()
                pres_id = pres_file["id"]

                self.logger.info(
                    f"✅ Uploaded presentation: {pres_id} - {safe_filename}"
                )

                ent = {
                    "type": "presentation",
                    "id": pres_id,
                    "filename": safe_filename,
                    "title": pres_content.title,
                    "token": token,
                    "expected_content": token,
                }
                out.append(ent)
                self._test_presentations.append(ent)
                self.created_entities.append({"id": pres_id, "name": safe_filename})

                self.logger.info(
                    f"📝 Presentation '{safe_filename}' created with token: {token}"
                )

        self.logger.info(
            f"✅ Created {len(self._test_presentations)} PowerPoint presentations"
        )
        return out

    def _sanitize_filename(self, filename: str) -> str:
        """Sanitize filename by removing illegal characters for OneDrive.

        Args:
            filename: Original filename

        Returns:
            Sanitized filename safe for OneDrive
        """
        # Replace illegal characters: \ / : * ? " < > |
        illegal_chars = ["\\", "/", ":", "*", "?", '"', "<", ">", "|"]
        safe_name = filename
        for char in illegal_chars:
            safe_name = safe_name.replace(char, "_")

        # Remove leading/trailing spaces and dots
        safe_name = safe_name.strip(". ")

        # Limit length to 200 characters (OneDrive has a 400 char limit for full path)
        if len(safe_name) > 200:
            # Keep the extension
            name, ext = (
                safe_name.rsplit(".", 1) if "." in safe_name else (safe_name, "")
            )
            safe_name = name[:195] + "." + ext if ext else name[:200]

        return safe_name

    def _create_powerpoint_file(self, pres_content: Any) -> bytes:
        """Create a PowerPoint presentation file with the given content.

        Args:
            pres_content: PowerPointPresentationContent object

        Returns:
            Bytes of the PowerPoint presentation
        """
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)

        # Slide 1: Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        title.text = pres_content.title
        subtitle.text = pres_content.subtitle

        # Add content slides
        for slide_content in pres_content.slides:
            # Use bullet layout
            bullet_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(bullet_slide_layout)

            # Set title
            title_shape = slide.shapes.title
            title_shape.text = slide_content.title

            # Add content
            content_shape = slide.placeholders[1]
            text_frame = content_shape.text_frame
            text_frame.clear()

            for i, bullet_text in enumerate(slide_content.content):
                if i == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                p.text = bullet_text
                p.level = 0

        # Save to bytes
        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        return buffer.getvalue()

    async def update_entities(self) -> List[Dict[str, Any]]:
        """Update PowerPoint presentations by adding a new slide with same tokens."""
        if not self._test_presentations:
            return []

        self.logger.info(
            f"🥁 Updating {min(2, len(self._test_presentations))} PowerPoint presentations"
        )
        updated = []

        async with httpx.AsyncClient(base_url=GRAPH, timeout=60) as client:
            for ent in self._test_presentations[
                : min(2, len(self._test_presentations))
            ]:
                await self._pace()

                # Download current presentation
                download_url = f"/me/drive/items/{ent['id']}/content"
                r = await client.get(download_url, headers=self._hdrs())

                if r.status_code != 200:
                    self.logger.warning(
                        f"Failed to download presentation: {r.status_code} - {r.text[:200]}"
                    )
                    continue

                # Load existing presentation
                prs = Presentation(io.BytesIO(r.content))

                # Add new slide with update information
                bullet_slide_layout = prs.slide_layouts[1]
                new_slide = prs.slides.add_slide(bullet_slide_layout)

                # Set title
                title_shape = new_slide.shapes.title
                title_shape.text = "Update Section"

                # Add content with token
                content_shape = new_slide.placeholders[1]
                text_frame = content_shape.text_frame
                text_frame.clear()

                update_content = [
                    "This presentation has been updated",
                    f"Update token: {ent['token']}",
                    "Updated to verify incremental sync functionality",
                ]

                for i, bullet_text in enumerate(update_content):
                    if i == 0:
                        p = text_frame.paragraphs[0]
                    else:
                        p = text_frame.add_paragraph()
                    p.text = bullet_text
                    p.level = 0

                # Save updated presentation
                buffer = io.BytesIO()
                prs.save(buffer)
                buffer.seek(0)
                updated_bytes = buffer.getvalue()

                # Upload updated presentation
                upload_url = f"/me/drive/items/{ent['id']}/content"
                r = await client.put(
                    upload_url,
                    headers={
                        "Authorization": f"Bearer {self.access_token}",
                        "Content-Type": (
                            "application/vnd.openxmlformats-officedocument."
                            "presentationml.presentation"
                        ),
                    },
                    content=updated_bytes,
                )

                if r.status_code in (200, 201):
                    updated.append({**ent, "updated": True})
                    self.logger.info(
                        f"📝 Updated presentation '{ent['filename']}' with token: {ent['token']}"
                    )
                else:
                    self.logger.warning(
                        f"Failed to update presentation: {r.status_code} - {r.text[:200]}"
                    )

                # Brief delay between updates
                await asyncio.sleep(0.5)

        return updated

    async def delete_entities(self) -> List[str]:
        """Delete all test presentations."""
        return await self.delete_specific_entities(self._test_presentations)

    async def delete_specific_entities(
        self, entities: List[Dict[str, Any]]
    ) -> List[str]:
        """Delete specific PowerPoint presentations with retry for locked files."""
        if not entities:
            # If no specific entities provided, delete all tracked presentations
            entities = self._test_presentations

        if not entities:
            return []

        self.logger.info(f"🥁 Deleting {len(entities)} PowerPoint presentations")
        deleted: List[str] = []

        async with httpx.AsyncClient(base_url=GRAPH, timeout=30) as client:
            for ent in entities:
                try:
                    await self._pace()

                    # Try to delete the presentation with retry for locked files (423)
                    max_retries = 3
                    retry_delay = 2.0  # seconds

                    for attempt in range(max_retries):
                        r = await client.delete(
                            f"/me/drive/items/{ent['id']}", headers=self._hdrs()
                        )

                        if r.status_code == 204:
                            deleted.append(ent["id"])
                            self.logger.info(
                                f"✅ Deleted presentation: {ent.get('filename', ent['id'])}"
                            )

                            # Remove from tracking
                            if ent in self._test_presentations:
                                self._test_presentations.remove(ent)
                            break  # Success, exit retry loop

                        elif r.status_code == 423 and attempt < max_retries - 1:
                            # Resource is locked, wait and retry
                            self.logger.warning(
                                "⏳ Presentation locked (423), retrying in %ss "
                                "(attempt %s/%s): %s",
                                retry_delay,
                                attempt + 1,
                                max_retries,
                                ent.get("filename", ent["id"]),
                            )
                            await asyncio.sleep(retry_delay)
                            retry_delay *= 2  # Exponential backoff
                        else:
                            # Other error or max retries reached
                            self.logger.warning(
                                f"Delete failed: {r.status_code} - {r.text[:200]}"
                            )
                            break  # Exit retry loop on non-retryable error

                except Exception as e:
                    self.logger.warning(
                        f"Delete error for {ent.get('filename', ent['id'])}: {e}"
                    )

        return deleted

    async def cleanup(self):
        """Comprehensive cleanup of all test resources."""
        self.logger.info("🧹 Starting comprehensive PowerPoint cleanup")

        cleanup_stats = {
            "presentations_deleted": 0,
            "errors": 0,
        }

        try:
            async with httpx.AsyncClient(base_url=GRAPH, timeout=30) as client:
                # Delete tracked test presentations
                if self._test_presentations:
                    self.logger.info(
                        f"🗑️  Deleting {len(self._test_presentations)} tracked presentations"
                    )
                    deleted = await self.delete_specific_entities(
                        self._test_presentations.copy()
                    )
                    cleanup_stats["presentations_deleted"] += len(deleted)

                # Search for and cleanup any orphaned test presentations
                await self._cleanup_orphaned_presentations(client, cleanup_stats)

            self.logger.info(
                f"🧹 Cleanup completed: {cleanup_stats['presentations_deleted']} "
                f"presentations deleted, {cleanup_stats['errors']} errors"
            )
        except Exception as e:
            self.logger.error(f"❌ Error during comprehensive cleanup: {e}")

    async def _cleanup_orphaned_presentations(
        self, client: httpx.AsyncClient, stats: Dict[str, Any]
    ):
        """Find and delete orphaned test presentations from previous runs."""
        try:
            await self._pace()
            r = await client.get("/me/drive/root/children", headers=self._hdrs())

            if r.status_code == 200:
                files = r.json().get("value", [])

                # Find test PowerPoint presentations
                test_presentations = [
                    f
                    for f in files
                    if f.get("name", "").startswith("Monke_")
                    and f.get("name", "").endswith(".pptx")
                ]

                if test_presentations:
                    self.logger.info(
                        f"🔍 Found {len(test_presentations)} orphaned test presentations"
                    )
                    for pres in test_presentations:
                        try:
                            await self._pace()
                            del_r = await client.delete(
                                f"/me/drive/items/{pres['id']}",
                                headers=self._hdrs(),
                            )
                            if del_r.status_code == 204:
                                stats["presentations_deleted"] += 1
                                self.logger.info(
                                    f"✅ Deleted orphaned presentation: "
                                    f"{pres.get('name', 'Unknown')}"
                                )
                            else:
                                stats["errors"] += 1
                        except Exception as e:
                            stats["errors"] += 1
                            self.logger.warning(
                                f"⚠️  Failed to delete presentation {pres['id']}: {e}"
                            )
        except Exception as e:
            self.logger.warning(f"⚠️  Could not search for orphaned presentations: {e}")

    def _hdrs(self) -> Dict[str, str]:
        """Get standard headers for Graph API requests."""
        return {
            "Authorization": f"Bearer {self.access_token}",
            "Content-Type": "application/json",
        }

    async def _pace(self):
        """Rate limiting helper."""
        now = time.time()
        if (delta := now - self._last_req) < self.rate_limit_delay:
            await asyncio.sleep(self.rate_limit_delay - delta)
        self._last_req = time.time()
